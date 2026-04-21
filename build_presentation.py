"""
Build SIPDiDS Semester 4 Team Project Presentation
Style: WSB lecture style — white slides, bold black titles, pink footer bar
Matches: Boolean Algebra L3.pdf reference style by Krzysztof Paszek
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ───────────────────────────────────────────────────────────────────
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x1A, 0x1A, 0x1A)
PINK     = RGBColor(0xE2, 0x00, 0x7A)   # WSB pink
GREY     = RGBColor(0x55, 0x55, 0x55)   # body text
LGREY    = RGBColor(0x88, 0x88, 0x88)   # sub-text / labels
DGREY    = RGBColor(0x33, 0x33, 0x33)   # dark heading
YELLOW   = RGBColor(0xF5, 0xC5, 0x18)   # title-slide gold bar

SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)
FOOTER_H = Inches(0.52)
FOOTER_Y = SLIDE_H - FOOTER_H

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def text(slide, txt, l, t, w, h,
         size=Pt(16), bold=False, color=GREY,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = txt
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_footer(slide, slide_num, section="SIPDIDS  |  TEAM DATA SCIENCE PROJECT — SEMESTER 4"):
    """Pink footer bar matching WSB lecture style."""
    rect(slide, 0, FOOTER_Y, SLIDE_W, FOOTER_H, fill=PINK)
    # WSB logo text (left)
    text(slide, "Akademia WSB", Inches(0.15), FOOTER_Y + Inches(0.04),
         Inches(2.5), Inches(0.25),
         size=Pt(10), bold=True, color=WHITE)
    text(slide, "WSB University", Inches(0.15), FOOTER_Y + Inches(0.26),
         Inches(2.5), Inches(0.2),
         size=Pt(9), bold=False, color=WHITE)
    # section label (centre)
    text(slide, section, Inches(2.8), FOOTER_Y + Inches(0.12),
         Inches(9.5), Inches(0.28),
         size=Pt(10), color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    # slide number (right)
    text(slide, str(slide_num), Inches(12.8), FOOTER_Y + Inches(0.12),
         Inches(0.45), Inches(0.28),
         size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

def white_bg(slide):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

def section_label(slide, label):
    """Small pink-text label above the title, like 'BOOLEAN ALGEBRA'."""
    text(slide, label, Inches(0.55), Inches(0.28), Inches(10), Inches(0.3),
         size=Pt(11), bold=True, color=PINK)

def slide_title(slide, title, y=Inches(0.62)):
    text(slide, title, Inches(0.55), y, Inches(12.2), Inches(1.0),
         size=Pt(34), bold=True, color=BLACK)

def pink_rule(slide, y=Inches(1.65)):
    rect(slide, Inches(0.55), y, Inches(12.2), Inches(0.025), fill=PINK)

def bullet_col(slide, items, l, t, w, h, item_h=None, color=GREY, size=Pt(14)):
    if item_h is None:
        item_h = h / max(len(items), 1)
    y = t
    for item in items:
        text(slide, f"• {item}", l, y, w, item_h + Inches(0.05),
             size=size, color=color)
        y += item_h

def box(slide, l, t, w, h, title, bullets,
        title_color=PINK, body_color=GREY,
        bg=RGBColor(0xF9,0xF9,0xF9), border=None):
    if border is None:
        border = title_color
    rect(slide, l, t, w, h, fill=bg, line_color=border, line_w=Pt(1.5))
    text(slide, title, l+Inches(0.15), t+Inches(0.08), w-Inches(0.3), Inches(0.38),
         size=Pt(13), bold=True, color=title_color)
    rect(slide, l, t+Inches(0.46), w, Inches(0.015), fill=border)
    y = t + Inches(0.52)
    bh = (h - Inches(0.6)) / max(len(bullets), 1)
    for b in bullets:
        text(slide, f"  {b}", l+Inches(0.1), y, w-Inches(0.2), bh+Inches(0.05),
             size=Pt(11), color=GREY)
        y += bh


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title (dark cover like the lecture template)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0x0D, 0x0D, 0x0D))

# WSB logo box (top centre)
rect(s, Inches(4.9), Inches(0.35), Inches(3.5), Inches(1.4),
     fill=WHITE, line_color=WHITE)
text(s, "Akademia WSB", Inches(5.05), Inches(0.45), Inches(3.2), Inches(0.45),
     size=Pt(18), bold=True, color=PINK, align=PP_ALIGN.CENTER)
text(s, "Dąbrowa Górnicza, Kraków, Cieszyn, Żywiec, Olkusz, Gliwice, Tychy",
     Inches(5.0), Inches(0.88), Inches(3.3), Inches(0.22),
     size=Pt(6), color=GREY, align=PP_ALIGN.CENTER)
rect(s, Inches(5.1), Inches(1.08), Inches(3.1), Inches(0.025), fill=PINK)
text(s, "WSB University", Inches(5.05), Inches(1.14), Inches(3.2), Inches(0.45),
     size=Pt(18), bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Main title
text(s, "SOCIAL INTELLIGENCE &\nPREDICTIVE DATA SCIENCE",
     Inches(0.5), Inches(2.35), Inches(12.3), Inches(1.8),
     size=Pt(32), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
text(s, "Team Data Science Project — Semester 4",
     Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
     size=Pt(16), color=RGBColor(0xCC,0xCC,0xCC), align=PP_ALIGN.CENTER)

# Gold rule
rect(s, Inches(5.4), Inches(4.85), Inches(2.5), Inches(0.055), fill=YELLOW)

# Team
text(s, "Percival  •  Shamil  •  Peris  •  Nihat  •  Christian",
     Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.4),
     size=Pt(13), color=RGBColor(0xBB,0xBB,0xBB), align=PP_ALIGN.CENTER,
     italic=True)

text(s, "Akademia WSB  |  April 2026",
     Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.35),
     size=Pt(11), color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

# Pink footer
rect(s, 0, FOOTER_Y, SLIDE_W, FOOTER_H, fill=PINK)
text(s, "Akademia WSB", Inches(0.15), FOOTER_Y+Inches(0.04),
     Inches(2.5), Inches(0.25), size=Pt(10), bold=True, color=WHITE)
text(s, "WSB University", Inches(0.15), FOOTER_Y+Inches(0.26),
     Inches(2.5), Inches(0.2), size=Pt(9), color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Section divider: Project Overview
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
text(s, "PROJECT OVERVIEW",
     Inches(2.5), Inches(2.8), Inches(8.3), Inches(1.1),
     size=Pt(46), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
text(s, "Lecture 2",
     Inches(2.5), Inches(3.95), Inches(8.3), Inches(0.5),
     size=Pt(18), color=LGREY, align=PP_ALIGN.CENTER, italic=True)
add_footer(s, 2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Overview content
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "PROJECT OVERVIEW")
slide_title(s, "Three Interconnected Research Streams")
pink_rule(s)
add_footer(s, 3)

LCARD = RGBColor(0xFD,0xF0,0xF7)
TCARD = RGBColor(0xF0,0xFB,0xF8)
PCARD = RGBColor(0xFD,0xF5,0xF9)

col_data = [
    (Inches(0.5),  SALMON:=RGBColor(0xFF,0x8C,0x69), LCARD, "Integration 1",
     "Percival + Shamil",
     ["Reddit social sentiment (PRAW API)",
      "Stock price & volatility (yfinance)",
      "5 tickers — 12-month window",
      "Pearson / Spearman correlation",
      "Q: Can Reddit predict stock moves?"]),
    (Inches(4.75), RGBColor(0x00,0x99,0x88), TCARD, "Integration 2",
     "Peris + Nihat",
     ["MovieLens dataset (25M ratings)",
      "VADER NLP sentiment on review text",
      "Welch t-test significance testing",
      "Shared movie dataset pipeline",
      "Q: Does sentiment match star rating?"]),
    (Inches(9.0),  PINK, PCARD, "Standalone",
     "Christian",
     ["CrossRef API — 5,500 records",
      "5 research fields × 10 years",
      "Publication volume & citation trends",
      "3-year moving average smoothing",
      "Q: Which fields are emerging/declining?"]),
]

for x, col, bg_col, tag, name, bullets in col_data:
    rect(s, x, Inches(1.82), Inches(4.1), Inches(4.38),
         fill=bg_col, line_color=col, line_w=Pt(2))
    text(s, tag, x+Inches(0.15), Inches(1.9), Inches(3.8), Inches(0.32),
         size=Pt(11), bold=True, color=col)
    text(s, name, x+Inches(0.15), Inches(2.22), Inches(3.8), Inches(0.38),
         size=Pt(16), bold=True, color=BLACK)
    rect(s, x+Inches(0.15), Inches(2.6), Inches(3.6), Inches(0.02), fill=col)
    y = Inches(2.7)
    for b in bullets:
        text(s, f"• {b}", x+Inches(0.2), y, Inches(3.7), Inches(0.54),
             size=Pt(12), color=GREY)
        y += Inches(0.54)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Section divider: Methodology
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
text(s, "METHODOLOGY",
     Inches(2.5), Inches(2.8), Inches(8.3), Inches(1.1),
     size=Pt(46), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
text(s, "The 7-Stage Data Mining Pipeline",
     Inches(2.5), Inches(3.95), Inches(8.3), Inches(0.5),
     size=Pt(18), color=LGREY, align=PP_ALIGN.CENTER, italic=True)
add_footer(s, 4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 7-Stage Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "METHODOLOGY")
slide_title(s, "The 7-Stage Data Mining Pipeline")
pink_rule(s)
add_footer(s, 5)

stages = [
    ("1", "Business\nUnderstanding", PINK),
    ("2", "Data\nAcquisition",       RGBColor(0xFF,0x8C,0x69)),
    ("3", "Pre-\nprocessing",        RGBColor(0x00,0xAA,0x88)),
    ("4", "Feature\nEngineering",    RGBColor(0x00,0xAA,0x88)),
    ("5", "Modelling",               RGBColor(0xFF,0x8C,0x69)),
    ("6", "Evaluation",              PINK),
    ("7", "Visualisation\n& Deploy", PINK),
]
bw = Inches(1.65)
bh = Inches(3.0)
sx = Inches(0.4)
sy = Inches(1.85)
gap = Inches(0.1)

for i, (num, label, col) in enumerate(stages):
    xx = sx + i*(bw+gap)
    bg_col = RGBColor(0xFF,0xF0,0xF7) if col == PINK else (
             RGBColor(0xFF,0xF5,0xEE) if col == RGBColor(0xFF,0x8C,0x69)
             else RGBColor(0xF0,0xFB,0xF8))
    rect(s, xx, sy, bw, bh, fill=bg_col, line_color=col, line_w=Pt(1.5))
    text(s, num, xx, sy+Inches(0.15), bw, Inches(0.65),
         size=Pt(30), bold=True, color=col, align=PP_ALIGN.CENTER)
    text(s, label, xx, sy+Inches(0.85), bw, Inches(1.8),
         size=Pt(12), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    if i < len(stages)-1:
        text(s, "→", xx+bw, sy+Inches(1.2), gap+Inches(0.05), Inches(0.5),
             size=Pt(18), color=LGREY, align=PP_ALIGN.CENTER)

text(s, "Applied consistently across all three research tasks  —  CRISP-DM inspired framework",
     Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.35),
     size=Pt(11), color=LGREY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Stages 1–2
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "STAGES 1 – 2")
slide_title(s, "Business Questions & Data Sources")
pink_rule(s)
add_footer(s, 6)

q_data = [
    (Inches(0.5),  PINK,                       "Integration 1  |  Percival + Shamil",
     ["Can Reddit sentiment predict stock price movement?",
      "Which companies show the strongest correlation?",
      "Do Reddit posts lead or lag stock events?",
      "Data: PRAW API + yfinance  (AMZN, TSLA, MSFT, GOOGL, NVDA)",
      "Period: 12 months of daily data"]),
    (Inches(4.75), RGBColor(0x00,0x99,0x88),   "Integration 2  |  Peris + Nihat",
     ["Does review sentiment align with numerical ratings?",
      "Which films show the largest sentiment–rating gap?",
      "Do high-rated films always get positive text reviews?",
      "Data: MovieLens 25M + review text corpus",
      "NLP tool: VADER compound score"]),
    (Inches(9.0),  RGBColor(0xFF,0x8C,0x69),   "Standalone  |  Christian",
     ["Which research fields grew most 2015–2025?",
      "Which fields are declining in publication volume?",
      "What drives citation density differences?",
      "Data: CrossRef API  (5 fields × 10 years)",
      "Records collected: 5,500 publication metadata entries"]),
]

for x, col, title, bullets in q_data:
    bg_col = RGBColor(0xFD,0xF0,0xF7) if col==PINK else (
             RGBColor(0xF0,0xFB,0xF8) if col==RGBColor(0x00,0x99,0x88)
             else RGBColor(0xFF,0xF5,0xEE))
    rect(s, x, Inches(1.82), Inches(4.1), Inches(4.3),
         fill=bg_col, line_color=col, line_w=Pt(1.5))
    text(s, title, x+Inches(0.15), Inches(1.92), Inches(3.8), Inches(0.38),
         size=Pt(12), bold=True, color=col)
    rect(s, x+Inches(0.15), Inches(2.3), Inches(3.6), Inches(0.02), fill=col)
    y = Inches(2.38)
    for b in bullets:
        text(s, f"• {b}", x+Inches(0.2), y, Inches(3.7), Inches(0.54),
             size=Pt(11.5), color=GREY)
        y += Inches(0.54)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Stages 3–4
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "STAGES 3 – 4")
slide_title(s, "Preprocessing & Feature Engineering")
pink_rule(s)
add_footer(s, 7)

box(s, Inches(0.5), Inches(1.82), Inches(6.0), Inches(3.55),
    "Data Cleaning — All Tasks",
    ["Removed null values, duplicates, and outliers",
     "Normalised date/time fields to consistent format",
     "Filtered to target time window (12 months / 10 years)",
     "Merged datasets on shared keys (ticker, movie_id)",
     "Verified record counts post-clean against raw totals"],
    title_color=PINK, bg=RGBColor(0xFD,0xF0,0xF7))

box(s, Inches(6.8), Inches(1.82), Inches(6.0), Inches(1.65),
    "Integration 1 — Feature Engineering",
    ["Lag features: Reddit posts D-1 to D-7 vs price D",
     "Rolling 7-day sentiment score per ticker",
     "Volatility: rolling std of daily returns"],
    title_color=RGBColor(0xFF,0x8C,0x69), bg=RGBColor(0xFF,0xF5,0xEE))

box(s, Inches(6.8), Inches(3.58), Inches(6.0), Inches(1.8),
    "Integration 2 — Feature Engineering",
    ["VADER compound → Positive / Neutral / Negative class",
     "Per-movie: average sentiment vs average star rating",
     "Discrepancy score = |sentiment_class − rating_quartile|",
     "Joined on movie_id (Peris → Nihat pipeline)"],
    title_color=RGBColor(0x00,0x99,0x88), bg=RGBColor(0xF0,0xFB,0xF8))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Section divider: Integration 1
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
text(s, "INTEGRATION 1",
     Inches(2.5), Inches(2.5), Inches(8.3), Inches(1.0),
     size=Pt(46), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
text(s, "Social Media × Stock Market",
     Inches(2.5), Inches(3.55), Inches(8.3), Inches(0.55),
     size=Pt(22), color=PINK, bold=True, align=PP_ALIGN.CENTER)
text(s, "Percival + Shamil",
     Inches(2.5), Inches(4.15), Inches(8.3), Inches(0.45),
     size=Pt(16), color=LGREY, align=PP_ALIGN.CENTER, italic=True)
add_footer(s, 8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Integration 1 Methodology
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "INTEGRATION 1  |  PERCIVAL + SHAMIL")
slide_title(s, "Social Media × Stock Market — Methodology")
pink_rule(s)
add_footer(s, 9)

box(s, Inches(0.5), Inches(1.82), Inches(6.1), Inches(3.55),
    "Percival — Reddit Social Media Analysis",
    ["PRAW API: scraped r/stocks, r/investing, r/wallstreetbets",
     "5 tickers: AMZN, TSLA, MSFT, GOOGL, NVDA",
     "Metrics: post volume, engagement score, sentiment polarity",
     "VADER NLP applied to post titles and body text",
     "Daily aggregated sentiment scores over 12 months"],
    title_color=PINK, bg=RGBColor(0xFD,0xF0,0xF7))

box(s, Inches(6.9), Inches(1.82), Inches(6.0), Inches(3.55),
    "Shamil — Stock Market Data Mining",
    ["yfinance API: OHLCV data, 12-month window",
     "Calculated: daily returns, 7-day / 30-day moving averages",
     "Volatility: rolling standard deviation of returns",
     "Lag correlation: Reddit signal D-1 to D-7 vs price D",
     "Pearson + Spearman correlation matrices per ticker"],
    title_color=RGBColor(0xFF,0x8C,0x69), bg=RGBColor(0xFF,0xF5,0xEE))

text(s, "Combined output: correlation heatmap  •  time-series overlay  •  lag feature importance chart",
     Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.3),
     size=Pt(11), color=LGREY, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Integration 1 Findings
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "INTEGRATION 1  |  KEY FINDINGS")
slide_title(s, "Reddit Sentiment vs Stock Price Movement")
pink_rule(s)
add_footer(s, 10)

findings = [
    (PINK,                       "Weak Overall Correlation",
     "Pearson r = 0.12 – 0.27 across tickers.\nReddit sentiment is a lagging indicator, not a leading one."),
    (RGBColor(0xFF,0x8C,0x69),   "AMZN Outlier",
     "Strongest link: AMZN at r = 0.27 (p = 0.03).\nHigh-volume posts preceded 3-day price upticks."),
    (RGBColor(0x00,0x99,0x88),   "Lag Effect Confirmed",
     "D-2 sentiment showed highest predictive value.\nImmediate-day (D-0) correlation was near zero."),
    (RGBColor(0x88,0x44,0xAA),   "Volatility Link",
     "High Reddit activity correlated with elevated volatility.\nNVDA & TSLA showed clearest spike–volatility pattern."),
]
x = Inches(0.5)
for col, title, body in findings:
    bg_col = RGBColor(0xFD,0xF0,0xF7)
    rect(s, x, Inches(1.82), Inches(2.95), Inches(3.6),
         fill=bg_col, line_color=col, line_w=Pt(2))
    text(s, title, x+Inches(0.15), Inches(1.9), Inches(2.65), Inches(0.45),
         size=Pt(13), bold=True, color=col)
    rect(s, x+Inches(0.15), Inches(2.35), Inches(2.65), Inches(0.02), fill=col)
    text(s, body, x+Inches(0.15), Inches(2.45), Inches(2.65), Inches(2.7),
         size=Pt(12), color=GREY)
    x += Inches(3.2)

text(s, "Visualisations: correlation_heatmap.png  •  timeseries_sentiment_price.png  •  viz_lag_feature_importance.png",
     Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.3),
     size=Pt(10), color=LGREY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Section divider: Integration 2
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
text(s, "INTEGRATION 2",
     Inches(2.5), Inches(2.5), Inches(8.3), Inches(1.0),
     size=Pt(46), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
text(s, "Movie Insights × Review Sentiment",
     Inches(2.5), Inches(3.55), Inches(8.3), Inches(0.55),
     size=Pt(22), color=RGBColor(0x00,0x99,0x88), bold=True, align=PP_ALIGN.CENTER)
text(s, "Peris + Nihat",
     Inches(2.5), Inches(4.15), Inches(8.3), Inches(0.45),
     size=Pt(16), color=LGREY, align=PP_ALIGN.CENTER, italic=True)
add_footer(s, 11)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Integration 2 Methodology
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "INTEGRATION 2  |  PERIS + NIHAT")
slide_title(s, "Movie Insights × Review Sentiment — Methodology")
pink_rule(s)
add_footer(s, 12)

box(s, Inches(0.5), Inches(1.82), Inches(6.1), Inches(3.55),
    "Peris — Movie Recommendation Insights",
    ["MovieLens: genre, runtime, release year, ratings",
     "Collected movie metadata via TMDB API",
     "Cleaned & structured: 1,200+ films with full attributes",
     "Feature correlation: avg_rating, num_ratings, genre_diversity",
     "Shared cleaned dataset with Nihat for sentiment join"],
    title_color=RGBColor(0x00,0x99,0x88), bg=RGBColor(0xF0,0xFB,0xF8))

box(s, Inches(6.9), Inches(1.82), Inches(6.0), Inches(3.55),
    "Nihat — Online Review Sentiment Analysis",
    ["VADER NLP: compound score on review text",
     "Classes: Positive (≥ 0.05) / Neutral / Negative",
     "Joined with Peris dataset on movie_id",
     "Per-movie: average sentiment vs average star rating",
     "Welch t-test: significant diff Positive vs Negative groups?"],
    title_color=PINK, bg=RGBColor(0xFD,0xF0,0xF7))

text(s, "Combined output: sentiment distribution  •  scatter sentiment vs rating  •  discrepancy highlight chart",
     Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.3),
     size=Pt(11), color=LGREY, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Integration 2 Findings
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "INTEGRATION 2  |  KEY FINDINGS")
slide_title(s, "Does Sentiment Match the Star Rating?")
pink_rule(s)
add_footer(s, 13)

findings2 = [
    (RGBColor(0x00,0x99,0x88), "Sentiment Distribution",
     "~60% positive  •  ~30% neutral  •  ~10% negative.\nReviews skew positive overall across the dataset."),
    (PINK,                     "Shawshank Effect",
     "High-rated films (4–5★) had disproportionately positive text.\nCult classics showed strongest sentiment–rating alignment."),
    (RGBColor(0xFF,0x8C,0x69), "Discrepancy Films",
     "Some films: high stars but neutral/negative text tone.\nSuggests polarising critical vs audience perception gap."),
    (RGBColor(0x44,0x88,0xCC), "Statistical Test",
     "Welch t-test: p < 0.05 — statistically significant.\nPositive reviews DO average higher numerical ratings."),
]
x = Inches(0.5)
for col, title, body in findings2:
    rect(s, x, Inches(1.82), Inches(2.95), Inches(3.6),
         fill=RGBColor(0xF8,0xF8,0xF8), line_color=col, line_w=Pt(2))
    text(s, title, x+Inches(0.15), Inches(1.9), Inches(2.65), Inches(0.45),
         size=Pt(13), bold=True, color=col)
    rect(s, x+Inches(0.15), Inches(2.35), Inches(2.65), Inches(0.02), fill=col)
    text(s, body, x+Inches(0.15), Inches(2.45), Inches(2.65), Inches(2.7),
         size=Pt(12), color=GREY)
    x += Inches(3.2)

text(s, "Visualisations: fig1_sentiment_distribution.png  •  fig2_scatter.png  •  fig4_discrepancy_highlight.png",
     Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.3),
     size=Pt(10), color=LGREY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Section divider: Standalone
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
text(s, "STANDALONE",
     Inches(2.5), Inches(2.5), Inches(8.3), Inches(1.0),
     size=Pt(46), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
text(s, "Academic Publication Mining",
     Inches(2.5), Inches(3.55), Inches(8.3), Inches(0.55),
     size=Pt(22), color=RGBColor(0xFF,0x8C,0x69), bold=True, align=PP_ALIGN.CENTER)
text(s, "Christian",
     Inches(2.5), Inches(4.15), Inches(8.3), Inches(0.45),
     size=Pt(16), color=LGREY, align=PP_ALIGN.CENTER, italic=True)
add_footer(s, 14)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Standalone Methodology
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "STANDALONE  |  CHRISTIAN")
slide_title(s, "Academic Publication Mining — Methodology")
pink_rule(s)
add_footer(s, 15)

box(s, Inches(0.5), Inches(1.82), Inches(6.0), Inches(3.55),
    "Data Acquisition",
    ["CrossRef API: queried by research field + year",
     "5 fields: AI/ML, Renewable Energy, Blockchain,",
     "  Bioinformatics, Quantum Computing",
     "Time range: 2015–2025  (10 years)",
     "Records collected: 5,500 publication metadata entries",
     "Fields captured: title, year, citations, DOI, keywords"],
    title_color=RGBColor(0xFF,0x8C,0x69), bg=RGBColor(0xFF,0xF5,0xEE))

box(s, Inches(6.8), Inches(1.82), Inches(6.0), Inches(3.55),
    "Analysis Pipeline",
    ["Publication volume trend per field per year",
     "3-year moving average to smooth year-to-year noise",
     "Year-over-year growth rate (%) calculation",
     "Citation density per paper by field",
     "Emerging vs declining: growth > +15% or < -10%",
     "Heatmap: field × year publication volume matrix"],
    title_color=PINK, bg=RGBColor(0xFD,0xF0,0xF7))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Standalone Findings
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "STANDALONE  |  KEY FINDINGS")
slide_title(s, "10 Years of Research Trends (2015–2025)")
pink_rule(s)
add_footer(s, 16)

box(s, Inches(0.5), Inches(1.82), Inches(3.9), Inches(3.55),
    "Emerging Fields",
    ["AI / ML — fastest absolute volume growth",
     "Renewable Energy — highest citation density per paper",
     "Quantum Computing — sharp post-2020 spike",
     "All three fields show consistent YoY growth"],
    title_color=RGBColor(0x00,0x99,0x88), bg=RGBColor(0xF0,0xFB,0xF8))

box(s, Inches(4.7), Inches(1.82), Inches(3.9), Inches(3.55),
    "Stable / Declining",
    ["Blockchain — peaked 2019–2021, now declining",
     "Bioinformatics — steady but growth rate slowing",
     "Citation density rising across all fields",
     "(depth of research increasing even if volume slows)"],
    title_color=RGBColor(0xFF,0x8C,0x69), bg=RGBColor(0xFF,0xF5,0xEE))

box(s, Inches(8.9), Inches(1.82), Inches(3.9), Inches(3.55),
    "Key Insight",
    ["Renewable Energy had the highest avg citations",
     "per paper across all 10 years",
     "→ High-impact, well-cited research even with",
     "  moderate publication volume growth",
     "→ Quality over quantity signal for this field"],
    title_color=PINK, bg=RGBColor(0xFD,0xF0,0xF7))

text(s, "Visualisations: Publication Volume heatmap  •  3-Year Moving Average  •  Emerging vs Declining trend chart",
     Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.3),
     size=Pt(10), color=LGREY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Evaluation
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "STAGE 6  |  EVALUATION")
slide_title(s, "How We Validated the Results")
pink_rule(s)
add_footer(s, 17)

box(s, Inches(0.5), Inches(1.82), Inches(3.9), Inches(3.55),
    "Integration 1 — Correlation Tests",
    ["Pearson r: linear correlation (normal dist.)",
     "Spearman ρ: rank correlation (non-parametric)",
     "Significance threshold: p < 0.05",
     "Result: AMZN r = 0.27, p = 0.03 ✓ significant",
     "Others: weak or non-significant (p > 0.05)"],
    title_color=PINK, bg=RGBColor(0xFD,0xF0,0xF7))

box(s, Inches(4.7), Inches(1.82), Inches(3.9), Inches(3.55),
    "Integration 2 — Welch T-Test",
    ["Compared avg ratings: Positive vs Negative groups",
     "Welch t-test: chosen for unequal group variance",
     "t-statistic and p-value computed per film cluster",
     "Result: p < 0.05 — sentiment class IS significant",
     "Effect size moderate — sentiment ≠ full predictor"],
    title_color=RGBColor(0x00,0x99,0x88), bg=RGBColor(0xF0,0xFB,0xF8))

box(s, Inches(8.9), Inches(1.82), Inches(3.9), Inches(3.55),
    "Standalone — Trend Validation",
    ["3-year moving average reduces volume noise",
     "YoY growth rate validates trend direction",
     "Citation density cross-validates field importance",
     "Descriptive trend analysis (no test needed)",
     "CrossRef records verified vs. known counts"],
    title_color=RGBColor(0xFF,0x8C,0x69), bg=RGBColor(0xFF,0xF5,0xEE))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Deployment / Dashboard
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "STAGE 7  |  DEPLOYMENT")
slide_title(s, "Integrated Streamlit Dashboard")
pink_rule(s)
add_footer(s, 18)

box(s, Inches(0.5), Inches(1.82), Inches(7.6), Inches(3.55),
    "master_dashboard.py — Features",
    ["Multi-page Streamlit app combining all 3 research tasks",
     "Home page: project overview, team bios, combined KPIs",
     "Page 1: Social Media × Stock correlation explorer",
     "Page 2: Sentiment vs Rating scatter & distribution",
     "Page 3: Academic publication volume & trend charts",
     "All charts interactive (Plotly)  —  WSB Charcoal & Rose theme"],
    title_color=PINK, bg=RGBColor(0xFD,0xF0,0xF7))

box(s, Inches(8.4), Inches(1.82), Inches(4.5), Inches(3.55),
    "Run Locally",
    ["cd team-project-semester4",
     "pip install -r requirements.txt",
     "streamlit run master_dashboard.py",
     "→ Opens at http://localhost:8501",
     "",
     "Next step: deploy to",
     "Streamlit Community Cloud (free)"],
    title_color=RGBColor(0xFF,0x8C,0x69), bg=RGBColor(0xFF,0xF5,0xEE))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Trello / Project Management
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "PROJECT MANAGEMENT")
slide_title(s, "Agile Workflow — Trello Kanban Board")
pink_rule(s)
add_footer(s, 19)

text(s, "Board: trello.com/b/85HqDUq1/team-project",
     Inches(0.5), Inches(1.78), Inches(8), Inches(0.32),
     size=Pt(11), color=LGREY, italic=True)

kanban = [
    ("To Do",       RGBColor(0x77,0x77,0x77), RGBColor(0xF5,0xF5,0xF5),
     "Task backlog defined at project kick-off. All 5 members assigned tasks upfront."),
    ("In Progress", RGBColor(0xFF,0x8C,0x69), RGBColor(0xFF,0xF5,0xEE),
     "Active work cards with owner name + deadline attached per card."),
    ("Review",      RGBColor(0xFF,0xCC,0x00), RGBColor(0xFF,0xFC,0xE8),
     "Peer-checking stage before marking done. Screenshots of output attached."),
    ("Done",        RGBColor(0x00,0x99,0x88), RGBColor(0xF0,0xFB,0xF8),
     "Completed tasks. Visualisation outputs attached as card attachments."),
]
x = Inches(0.5)
for col_name, col, bg_col, desc in kanban:
    rect(s, x, Inches(2.2), Inches(3.0), Inches(3.1),
         fill=bg_col, line_color=col, line_w=Pt(2))
    rect(s, x, Inches(2.2), Inches(3.0), Inches(0.42), fill=col)
    text(s, col_name, x+Inches(0.1), Inches(2.23), Inches(2.8), Inches(0.38),
         size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, desc, x+Inches(0.15), Inches(2.72), Inches(2.7), Inches(2.45),
         size=Pt(12), color=GREY)
    x += Inches(3.2)

text(s, "All 5 team members tracked  •  Task deadlines enforced via Trello due dates  •  Screenshots attached to cards",
     Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.35),
     size=Pt(11), color=LGREY, italic=True, align=PP_ALIGN.CENTER)
text(s, "[Insert Trello screenshots from: ...\\Sem 4\\SIPDiDS Team Project\\trello\\  onto this slide]",
     Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.3),
     size=Pt(10), color=PINK, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusions
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)
section_label(s, "CONCLUSIONS")
slide_title(s, "What We Learned Across All Three Tasks")
pink_rule(s)
add_footer(s, 20)

conclusions = [
    (PINK,                       "Social signals lag, not lead",
     "Reddit sentiment weakly predicts stock movement with a 2-day lag — useful context, not a reliable standalone predictor."),
    (RGBColor(0x00,0x99,0x88),   "Text and numbers tell different stories",
     "Review sentiment aligns with ratings overall but significant discrepancies exist — NLP captures nuance that star ratings miss."),
    (RGBColor(0xFF,0x8C,0x69),   "Research fields have life cycles",
     "AI/ML and Renewable Energy are clearly ascending; Blockchain peaked and declined — publication data tracks paradigm shifts."),
    (RGBColor(0x44,0x88,0xCC),   "Shared pipelines reduce friction",
     "Cross-team data sharing (Peris→Nihat, Percival→Shamil) standardised preprocessing and accelerated the integration work."),
]

y = Inches(1.85)
for col, title, body in conclusions:
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.85),
         fill=RGBColor(0xF8,0xF8,0xF8), line_color=col, line_w=Pt(2))
    rect(s, Inches(0.5), y, Inches(0.08), Inches(0.85), fill=col)
    text(s, title, Inches(0.75), y+Inches(0.08), Inches(3.3), Inches(0.38),
         size=Pt(13), bold=True, color=col)
    text(s, body,  Inches(4.2),  y+Inches(0.08), Inches(8.4), Inches(0.72),
         size=Pt(12), color=GREY)
    y += Inches(1.0)

text(s, "All findings reproducible — code, data, and dashboard: github.com/Mahwaya/team-project-semester4",
     Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.3),
     size=Pt(10), color=LGREY, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Q&A
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_bg(s)

text(s, "Thank You",
     Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.4),
     size=Pt(58), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
text(s, "Questions & Discussion",
     Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.7),
     size=Pt(26), color=PINK, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(4.0), Inches(3.85), Inches(5.3), Inches(0.04), fill=PINK)

text(s, "Percival  •  Shamil  •  Peris  •  Nihat  •  Christian",
     Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.45),
     size=Pt(15), color=DGREY, align=PP_ALIGN.CENTER)

text(s, "github.com/Mahwaya/team-project-semester4",
     Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.4),
     size=Pt(14), bold=True, color=RGBColor(0x00,0x99,0x88), align=PP_ALIGN.CENTER)

text(s, "MSc Computer Science (Data Science)  |  Akademia WSB  |  Semester 4  |  April 2026",
     Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.35),
     size=Pt(11), color=LGREY, align=PP_ALIGN.CENTER)

add_footer(s, 21)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = r"C:\Users\Percival Mahwaya\Desktop\SIPDiDS_Team_Presentation_Sem4.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
